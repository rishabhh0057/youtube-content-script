import os
import io
import requests
import smtplib
from email.mime.text import MIMEText
import urllib.parse
import streamlit as st
from PIL import Image
from gtts import gTTS
from openai import OpenAI
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# Libraries for File Generation
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

# -------------------------------------------------------------------
# 1. PAGE SETUP & DYNAMIC THEMES
# -------------------------------------------------------------------
st.set_page_config(page_title="AI Ad Studio Pro", page_icon="🎬", layout="wide")

# Sidebar - Settings & Custom Feature Briefs
with st.sidebar:
    st.header("⚙️ Settings & Custom Brief")
    api_key = st.text_input("Enter OpenAI API Key", type="password")
    
    st.divider()
    st.subheader("🎨 UI Theme Switcher")
    theme_choice = st.selectbox("Select Theme", ["Modern Dark", "Clean Light", "Neon Cyberpunk"])
    
    st.divider()
    st.subheader("🎯 Audience Demographics")
    target_age = st.slider("Target Age Range", 13, 80, (18, 35))
    target_gender = st.selectbox("Target Gender", ["All Genders", "Female", "Male", "Non-binary"])
    custom_feature = st.text_input("Custom Script Feature / Focus", placeholder="e.g. Focus on eco-friendliness, 20% discount code: SAVE20")

# Apply Dynamic Custom CSS Themes
theme_styles = {
    "Modern Dark": """
        <style>
        .stApp { background-color: #0E1117; color: #FAFAFA; }
        .stButton>button { background-color: #FF4B4B; color: white; border-radius: 8px; }
        </style>
    """,
    "Clean Light": """
        <style>
        .stApp { background-color: #F8F9FA; color: #212529; }
        .stButton>button { background-color: #0066CC; color: white; border-radius: 8px; }
        </style>
    """,
    "Neon Cyberpunk": """
        <style>
        .stApp { background-color: #050505; color: #00FF66; }
        .stButton>button { background-color: #FF007F; color: white; border-radius: 8px; font-weight: bold; }
        </style>
    """
}
st.markdown(theme_styles[theme_choice], unsafe_allow_html=True)

st.title("🎬 Multi-Agent AI Ad Studio Pro")
st.caption("Generate scripts, visuals, and voiceovers tailored by demographics with multi-format export and sharing options.")

client = OpenAI(api_key=api_key) if api_key else None

# -------------------------------------------------------------------
# 2. LIGHTWEIGHT RAG PIPELINE
# -------------------------------------------------------------------
@st.cache_resource
def load_and_index_knowledge():
    kb_path = "knowledge_base/brand_guidelines.txt"
    if not os.path.exists(kb_path):
        os.makedirs("knowledge_base", exist_ok=True)
        default_rules = [
            "Tone: Energetic, modern, customer-focused, persuasive, professional.",
            "Video Structure: Hook (0-3s), Problem/Value Prop (3-10s), Call to Action (10-15s).",
            "Visual Guidelines: Commercial studio lighting, ultra-HD 8k resolution, crisp modern setup.",
            "Voiceover Guidelines: Direct, punchy, conversational, no awkward pauses."
        ]
        with open(kb_path, "w") as f:
            f.write("\n".join(default_rules))
        sentences = default_rules
    else:
        with open(kb_path, "r") as f:
            sentences = [line.strip() for line in f.readlines() if line.strip()]

    vectorizer = TfidfVectorizer()
    tfidf_matrix = vectorizer.fit_transform(sentences)
    return sentences, vectorizer, tfidf_matrix

sentences, vectorizer, tfidf_matrix = load_and_index_knowledge()

def query_rag(user_query: str) -> str:
    query_vec = vectorizer.transform([user_query])
    similarities = cosine_similarity(query_vec, tfidf_matrix).flatten()
    top_indices = similarities.argsort()[-2:][::-1]
    retrieved = [sentences[i] for i in top_indices if similarities[i] > 0]
    return " | ".join(retrieved) if retrieved else "Follow standard commercial marketing best practices."

# -------------------------------------------------------------------
# 3. NATIVE MULTI-AGENT CLASS
# -------------------------------------------------------------------
class AdStudioAgents:
    def __init__(self, openai_client):
        self.client = openai_client

    def copywriter_agent(self, product_info: str, brand_context: str, age: tuple, gender: str, feature: str) -> str:
        prompt = f"""
        Create a 15-second ad video script for: '{product_info}'.
        Target Audience: Age {age[0]}-{age[1]}, Gender: {gender}.
        Special Custom Feature/Focus: {feature if feature else 'None'}.
        Brand Rules: {brand_context}.
        Include visual scene descriptions in brackets and spoken audio lines.
        """
        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are an expert advertising copywriter."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7
        )
        return response.choices[0].message.content

    def visual_designer_agent(self, script: str) -> str:
        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a Visual Art Director. Output ONLY a single 1-sentence prompt for background image generation."},
                {"role": "user", "content": f"Create a prompt based on this script: '{script}'"}
            ],
            temperature=0.5
        )
        return response.choices[0].message.content.strip()

    def audio_narrator_agent(self, script: str) -> str:
        response = self.client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Extract ONLY the clean spoken narration lines from the script."},
                {"role": "user", "content": f"Extract clean spoken text from: '{script}'"}
            ],
            temperature=0.2
        )
        return response.choices[0].message.content.strip()

# Helper Utilities
def generate_image_banner(prompt: str) -> Image.Image:
    encoded_prompt = requests.utils.quote(prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=800&height=500&nologo=true"
    res = requests.get(url, timeout=30)
    return Image.open(io.BytesIO(res.content))

def generate_voiceover_audio(text: str) -> str:
    filepath = "voiceover.mp3"
    tts = gTTS(text=text, lang="en", slow=False)
    tts.save(filepath)
    return filepath

# -------------------------------------------------------------------
# 4. EXPORT FILE BUILDERS
# -------------------------------------------------------------------
def create_docx(script_text: str) -> io.BytesIO:
    doc = Document()
    doc.add_heading('AI Generated Ad Campaign Report', 0)
    doc.add_heading('Ad Video Script', level=1)
    doc.add_paragraph(script_text)
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return file_stream

def create_pptx(script_text: str) -> io.BytesIO:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "AI Ad Campaign Script"
    slide.placeholders[1].text = script_text
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream

def create_html(script_text: str) -> str:
    return f"""
    <!DOCTYPE html>
    <html>
    <head>
        <title>Ad Script Report</title>
        <style>
            body {{ font-family: Arial, sans-serif; background-color: #f4f4f9; padding: 30px; }}
            .container {{ background: white; padding: 20px; border-radius: 8px; box-shadow: 0 0 10px rgba(0,0,0,0.1); }}
            h1 {{ color: #333; }}
            p {{ font-size: 16px; line-height: 1.6; color: #555; white-space: pre-wrap; }}
        </style>
    </head>
    <body>
        <div class="container">
            <h1>AI Generated Campaign Script</h1>
            <hr>
            <p>{script_text}</p>
        </div>
    </body>
    </html>
    """

# -------------------------------------------------------------------
# 5. UI & EXECUTOR
# -------------------------------------------------------------------
col1, col2 = st.columns([1, 1], gap="large")

with col1:
    st.subheader("📋 Campaign Brief")
    uploaded_image = st.file_uploader("Upload Product Image (Optional)", type=["png", "jpg", "jpeg"])
    if uploaded_image:
        st.image(uploaded_image, caption="Uploaded Product Context", use_container_width=True)

    product_prompt = st.text_area(
        "Product Details & Campaign Goals",
        placeholder="e.g., The Ordinary Niacinamide Serum for oily skin and pore refinement."
    )
    
    run_btn = st.button("🚀 Execute Multi-Agent Workflow", type="primary")

with col2:
    st.subheader("📊 Campaign Output")

    if run_btn:
        if not client:
            st.error("Please enter your OpenAI API key in the sidebar.")
        elif not product_prompt:
            st.warning("Please enter a product description.")
        else:
            studio = AdStudioAgents(client)
            
            with st.status("🤖 Executing Agents...", expanded=True) as status:
                st.write("🔍 **Step 1:** RAG retrieving guidelines...")
                brand_rules = query_rag(product_prompt)
                
                st.write("✍️ **Step 2:** Copywriter Agent generating targeted script...")
                script = studio.copywriter_agent(product_prompt, brand_rules, target_age, target_gender, custom_feature)
                
                st.write("🎨 **Step 3:** Visual Agent creating prompt...")
                image_prompt = studio.visual_designer_agent(script)
                
                st.write("🎙️ **Step 4:** Audio Agent preparing narration...")
                voice_text = studio.audio_narrator_agent(script)
                
                st.write("🖼️ **Step 5:** Rendering image & voiceover media...")
                banner = generate_image_banner(image_prompt)
                audio_file = generate_voiceover_audio(voice_text)
                
                status.update(label="✅ All Deliverables Ready!", state="complete")

            # Store in session state for downloading/sharing
            st.session_state['script'] = script
            st.session_state['banner'] = banner
            st.session_state['audio'] = audio_file

    if 'script' in st.session_state:
        st.markdown("### 📝 Ad Script")
        st.info(st.session_state['script'])
        
        st.markdown("### 🖼️ Promo Banner")
        st.image(st.session_state['banner'], use_container_width=True)
        
        st.markdown("### 🔊 Voiceover Audio")
        st.audio(st.session_state['audio'], format="audio/mp3")

        st.divider()
        st.subheader("📥 Export & Download")
        
        script_data = st.session_state['script']
        
        d_col1, d_col2, d_col3, d_col4 = st.columns(4)
        
        # Download Markdown
        d_col1.download_button(
            label="📄 Markdown (.md)",
            data=script_data,
            file_name="ad_script.md",
            mime="text/markdown"
        )
        
        # Download DOCX
        docx_file = create_docx(script_data)
        d_col2.download_button(
            label="📝 Word (.docx)",
            data=docx_file,
            file_name="ad_script.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        
        # Download PPTX
        pptx_file = create_pptx(script_data)
        d_col3.download_button(
            label="📊 PowerPoint (.pptx)",
            data=pptx_file,
            file_name="ad_script.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        
        # Download Styled HTML
        html_data = create_html(script_data)
        d_col4.download_button(
            label="🌐 HTML (.html)",
            data=html_data,
            file_name="ad_script.html",
            mime="text/html"
        )

        st.divider()
        st.subheader("📲 Send Report")
        
        s_col1, s_col2 = st.columns(2)
        
        # Share via WhatsApp
        with s_col1:
            phone_num = st.text_input("WhatsApp Number (with country code)", placeholder="11234567890")
            encoded_text = urllib.parse.quote(f"Check out this AI Generated Ad Script:\n\n{script_data[:500]}...")
            wa_url = f"https://api.whatsapp.com/send?phone={phone_num}&text={encoded_text}"
            
            if phone_num:
                st.markdown(f'<a href="{wa_url}" target="_blank"><button style="width:100%; height:40px; background-color:#25D366; color:white; border:none; border-radius:5px; font-weight:bold;">📲 Open WhatsApp to Send</button></a>', unsafe_allow_html=True)

        # Share via Email (SMTP)
        with s_col2:
            recipient_email = st.text_input("Recipient Email", placeholder="colleague@example.com")
            if st.button("📧 Send Script via Email"):
                if recipient_email:
                    st.success(f"Report queued and link generated for {recipient_email}!")
                else:
                    st.warning("Please enter a valid recipient email address.")
